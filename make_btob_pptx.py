"""BtoB AX支援サービス素案を編集可能な.pptx(10枚)で生成。文字・表はPowerPoint/Googleスライドで直接編集可。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK=RGBColor(0x16,0x2a,0x4a); INK2=RGBColor(0x3c,0x54,0x78); GREEN=RGBColor(0x5c,0x90,0x5a)
GREEN_D=RGBColor(0x46,0x76,0x46); ACCENT=RGBColor(0xe4,0x8e,0x3c); BG=RGBColor(0xee,0xf3,0xe6)
LIGHT=RGBColor(0xf5,0xf8,0xf0); GRAY=RGBColor(0x60,0x6e,0x64); WHITE=RGBColor(0xff,0xff,0xff)
CARD=RGBColor(0xfc,0xfd,0xf9); GREEN2=RGBColor(0x96,0xbc,0x8c)
FONT="Yu Gothic"
OUT="C:/Users/takanori.koyama/Desktop/Claude/sns-assets/AX支援サービス素案_GREE.pptx"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide():
    s=prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb=BG
    return s

def box(s,l,t,w,h,fill=None,line=None,line_w=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(line_w)
    return sp

def txt(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp_after=4):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    if isinstance(runs,str): runs=[(runs,18,INK,False)]
    first=True
    for item in runs:
        text,size,color,bold=item
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align; p.space_after=Pt(sp_after)
        r=p.add_run(); r.text=text; f=r.font
        f.size=Pt(size); f.bold=bold; f.color.rgb=color; f.name=FONT
    return tb

def title(s,text,page):
    box(s,0.4,0.32,0.14,0.62,fill=GREEN)
    txt(s,0.62,0.30,11.4,0.8,[(text,30,INK,True)],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,12.2,0.34,0.9,0.5,[(f"{page} / 10",14,GRAY,False)])
    ln=box(s,0.4,1.12,12.53,0.02,fill=GREEN2)
    txt(s,0.4,7.02,6,0.4,[("GREE｜AX支援サービス 素案",11,GRAY,False)])

def card(s,l,t,w,h,header,color,lines=None,header_fs=18):
    box(s,l,t,w,h,fill=CARD,line=color,line_w=2.2)
    if header:
        box(s,l,t,w,0.62,fill=color,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s,l,t+0.04,w,0.54,[(header,header_fs,WHITE,True)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    return

def table(s,l,t,w,h,data,col_w,fs=11,hi_row=None):
    rows=len(data); cols=len(data[0])
    gt=s.shapes.add_table(rows,cols,Inches(l),Inches(t),Inches(w),Inches(h)).table
    gt.first_row=False; gt.horz_banding=False
    total=sum(col_w)
    for c,cw in enumerate(col_w): gt.columns[c].width=Emu(int(Inches(w)*cw/total))
    for r in range(rows):
        for c in range(cols):
            cell=gt.cell(r,c); cell.margin_left=Inches(0.08); cell.margin_right=Inches(0.06)
            cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
            cell.vertical_anchor=MSO_ANCHOR.MIDDLE
            if r==0: cell.fill.solid(); cell.fill.fore_color.rgb=GREEN
            elif hi_row is not None and r==hi_row: cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(0xff,0xf4,0xe4)
            elif r%2==0: cell.fill.solid(); cell.fill.fore_color.rgb=LIGHT
            else: cell.fill.solid(); cell.fill.fore_color.rgb=CARD
            tf=cell.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
            run=p.add_run(); run.text=str(data[r][c]); f=run.font
            f.size=Pt(fs); f.name=FONT
            f.bold=(r==0 or c==0)
            f.color.rgb=WHITE if r==0 else (ACCENT if (hi_row==r and c==0) else INK)
    return gt

# S1 表紙
s=slide(); box(s,0,0,13.333,0.16,fill=GREEN)
txt(s,1,2.0,11.33,1.2,[("BtoB AX支援サービス（素案）",46,INK,True)],align=PP_ALIGN.CENTER)
txt(s,1,3.2,11.33,0.7,[("SNS運用 × AI × 人 × SocialPitt",26,GREEN_D,True)],align=PP_ALIGN.CENTER)
txt(s,1,4.1,11.33,1.0,[("「続ける」だけでなく「伸ばす」。",22,INK2,False),
   ("フォロワー・エンゲージメント拡大まで、代理店が伴走する。",22,INK2,False)],align=PP_ALIGN.CENTER)
box(s,4.9,5.6,3.5,0.6,fill=GREEN); txt(s,4.9,5.62,3.5,0.56,[("GREE ／ 2026-06",16,WHITE,True)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# S2 課題
s=slide(); title(s,"よくある悩み：投稿は続いた。でも“伸びない”。",2)
card(s,0.5,1.5,5.9,3.4,"AIだけだと…",GRAY,None)
txt(s,0.8,2.4,5.3,2.2,[("効率（量・継続・一貫性）… ○ 出る",18,GREEN_D,True),
  ("成果（フォロワー・ENG）… △ 不透明",18,ACCENT,True),
  ("",10,INK,False),("量産はできても、“伸ばす”仕掛けが足りない。",16,INK2,False)])
card(s,6.9,1.5,5.9,3.4,"現場のペイン",ACCENT,None)
txt(s,7.2,2.45,5.3,2.2,[("・ネタ切れ・制作工数で続かない",18,INK,False),
  ("・トンマナ・品質がぶれる",18,INK,False),("・何を出せばいいか戦略がない",18,INK,False),
  ("・効果が見えない／改善できない",18,INK,False)])
txt(s,0.5,5.4,12.3,0.7,[("「効率」と「成果」は別物。成果には“もう2つの力”が要る。",26,ACCENT,True)],align=PP_ALIGN.CENTER)

# S3 3レイヤー
s=slide(); title(s,"解決：3つの力を組み合わせる",3)
data3=[("AI（仕組み）",GREEN,"続ける／そろえる","半自動で量産・キャラ統一・予約投稿"),
 ("人（GREE）",ACCENT,"当てる／磨く","戦略・編集・コメント運用・改善判断"),
 ("SocialPitt",GREEN_D,"測る／回す","KPI可視化・競合ベンチ・勝ち筋抽出")]
for i,(h,c,role,desc) in enumerate(data3):
    l=0.5+i*4.15; card(s,l,1.5,3.85,2.6,h,c)
    txt(s,l,2.35,3.85,0.6,[(role,22,c,True)],align=PP_ALIGN.CENTER)
    txt(s,l+0.25,3.05,3.35,1.0,[(desc,15,INK,False)],align=PP_ALIGN.CENTER)
txt(s,0.5,4.5,12.3,0.5,[("データ(SocialPitt) → 人が“次の企画”に翻訳 → AIが量産 → 投稿・計測（PDCA）",18,INK,True)],align=PP_ALIGN.CENTER)
txt(s,0.5,5.3,12.3,0.6,[("このループが回るほど伸びる",24,GREEN_D,True)],align=PP_ALIGN.CENTER)

# S4 KPI
s=slide(); title(s,"なぜ“伸びる”のか（フォロワーの増え方）",4)
box(s,0.5,1.4,12.3,1.3,fill=LIGHT,line=GREEN,line_w=2)
txt(s,0.5,1.55,12.3,1.1,[("フォロワー純増 ＝ リーチ × プロフィール訪問率 × フォロー転換率 −（離脱）",22,INK,True),
  ("ENG率 ＝ （保存＋いいね＋コメント＋シェア）÷ リーチ",18,INK2,False)],align=PP_ALIGN.CENTER)
table(s,0.5,3.0,12.3,3.0,[["伸ばす要素","担い手","具体"],
 ["露出量（量・継続・一貫性）","AI","半自動で量産・毎日投稿"],
 ["当たり化（実用性・フック・トンマナ）","人(GREE)","企画・編集・見せ方"],
 ["最適化（時間/タグ/競合比較）","SocialPitt","データで勝ち筋を発見"],
 ["関係づくり（コメント返信）","人(GREE)","ファン化→評価UP"]],[5,2,5],fs=14)

# S5 PDCA
s=slide(); title(s,"回し方（毎週まわすPDCA）",5)
cyc=[("①分析 (SocialPitt)",GREEN),("②企画 (人)",GREEN),("③量産 (AI)",GREEN),("④承認 (顧客)",ACCENT),("⑤投稿→運用(人)",GREEN)]
for i,(t,c) in enumerate(cyc):
    l=0.4+i*2.55; card(s,l,1.7,2.3,1.3,None,c)
    txt(s,l,1.7,2.3,1.3,[(t,15,(ACCENT if c==ACCENT else INK),True)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
txt(s,0.5,3.4,12.3,0.6,[("週次：勝ち筋を翌週の企画へ反映　／　月次：戦略・テーマ配分・広告ブーストを見直し",16,INK2,False)],align=PP_ALIGN.CENTER)
box(s,3.5,4.3,6.3,0.9,fill=LIGHT,line=ACCENT,line_w=2)
txt(s,3.5,4.3,6.3,0.9,[("④承認（顧客）＝“勝手に出さない”安全装置",20,ACCENT,True)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# S6 権限・セキュリティ
s=slide(); title(s,"安心して任せられる仕組み（権限・セキュリティ）",6)
table(s,0.5,1.4,12.3,4.6,[["観点","どうする"],
 ["アカウント所有","クライアントが保有（代理店に譲渡しない）"],
 ["代理店の権限","Metaパートナー連携で“最小権限”のみ付与"],
 ["代行投稿","可。ただし最終OK（承認）は必ず顧客"],
 ["トークン/秘密情報","顧客資産として管理・平文コミット禁止・漏洩時は即再発行"],
 ["記録","誰が承認・投稿したかの監査ログを保持"],
 ["AIに入れない","顧客の機密・個人情報は投入しない"],
 ["ブランドセーフティ","炎上回避チェック（誇大・不安訴求・権利侵害を排除）"]],[3,9],fs=13)
box(s,0.5,6.2,7.0,0.55,fill=ACCENT)
txt(s,0.5,6.2,7.0,0.55,[("キモ：代行はするが“最終OKは顧客”＋データは守る",16,WHITE,True)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# S7 パッケージ
s=slide(); title(s,"料金プラン（4段階・成果志向）",7)
table(s,0.3,1.35,12.7,4.4,[["プラン","狙い","AI","人(GREE)","SocialPitt","投稿/レポート","KPI"],
 ["① Starter","立ち上げ","設計・テンプレ・初期投稿","戦略設計・初期伴走","KPI設計＋ダッシュボード","月8〜12／月次","基盤・初期母数"],
 ["② Standard","続けて磨く","半自動量産・予約","編集監修・週次改善","週次分析・ベンチ","月12〜20／隔週","ENG率改善"],
 ["③ Growth ★","伸ばす","量産＋A/B＋リール","プランナー＋アナリスト","高度分析・競合/タグ最適化","月20〜30＋リール／週次","フォロワー純増・ENGに数値コミット"],
 ["④ Enablement","内製化/研修","仕組み移管・プロンプト資産","研修・伴走・レビュー","自社運用へ移管","顧客運用／月次","内製能力の獲得"]],
 [1.4,1.0,1.7,1.6,1.7,1.6,1.9],fs=11,hi_row=3)
txt(s,0.3,6.2,12.7,0.6,[("成果（フォロワー/ENG）が要件なら → 人＋SocialPittの比重が高い ③Growth が必須",18,ACCENT,True)],align=PP_ALIGN.CENTER)

# S8 座組マトリクス
s=slide(); title(s,"提供の座組（A／B／C）— AI活用の観点で",8)
table(s,0.3,1.35,12.7,4.7,[["観点","A：フル外注","B：部分外注","C：AI内製化コンサル"],
 ["クライアントの役割","承認＋レポート受領＋継続判断のみ","自社制作・投稿は内製／一部を選択外注","運用は自社（仕組みを受け取る）"],
 ["GREEの役割","企画〜制作〜投稿〜分析〜改善まで全部代行","選択領域(キャンペーン/UGC/SocialPitt)を支援","Claude活用の内製化を設計・実装・研修"],
 ["AI活用（誰が使う）","GREEがAIで量産・効率化（顧客はAIレス）","役割分担：外注=GREEのAI／内製=顧客","顧客自身がAIで作れるように（民主化）"],
 ["SocialPitt","GREEが運用しレポート","分析だけ単体外注も可","自社ダッシュボードとして導入支援"],
 ["人(GREE)の関与","高（フル運用）","中（選択領域のみ）","初期=高 → 移管後=低"],
 ["向く顧客","リソース無い／丸ごと任せたい","内製体制あり／繁忙・専門領域を補完","内製化・恒久コスト削減・ノウハウ蓄積"],
 ["課金イメージ","月額運用（高）","メニュー/PJ＋SocialPitt月額","コンサル/実装PJ＋研修"],
 ["対応パッケージ／現状","②〜③ / 提供中","①②＋オプション / 提供中","④Enablement / 新規"]],[2.4,3.4,3.4,3.4],fs=10.5)

# S9 AI活用3類型
s=slide(); title(s,"AI活用の3類型（座組の本質＝“誰がAIを操るか”）",9)
t3=[("代行型（A）",GREEN,"GREEがAIで“代わりに”回す","手離れ・継続・成果重視","②③ / フル外注"),
    ("ハイブリッド型（B）",ACCENT,"AIの使い手を分担する","柔軟・コスト最適・補完","① 部分＋SocialPitt"),
    ("内製化型（C）",GREEN_D,"顧客がAIを使えるように","自走・恒久効率・ノウハウ","④ Enablement")]
for i,(h,c,role,val,pkg) in enumerate(t3):
    l=0.5+i*4.15; card(s,l,1.45,3.85,3.4,h,c)
    txt(s,l,2.3,3.85,0.6,[(role,18,c,True)],align=PP_ALIGN.CENTER)
    txt(s,l+0.3,3.1,3.25,1.4,[("提供価値",13,GRAY,True),(val,16,INK,False),("",8,INK,False),
       ("対応",13,GRAY,True),(pkg,15,c,True)])
box(s,0.5,5.2,12.3,1.0,fill=LIGHT,line=GREEN_D,line_w=2)
txt(s,0.5,5.25,12.3,0.95,[("どれもAIが“てこ”。違いは「誰がAIを操るか」だけ。",24,GREEN_D,True),
  ("顧客の体制・目的に合わせて A/B/C を選ぶ（組み合わせも可）",16,INK2,False)],align=PP_ALIGN.CENTER)

# S10 まとめ
s=slide(); title(s,"はじめ方（3ステップ）とまとめ",10)
steps=[("STEP 1","自社で実証","ai_coach_korekara を実績化（100フォロワー＆運用データ）"),
 ("STEP 2","パイロット1社","Growthで試行→フロー・権限・セキュリティを実地検証"),
 ("STEP 3","横展開","標準パッケージ化→最後に提案資料を作成")]
for i,(st,t,desc) in enumerate(steps):
    l=0.5+i*4.15; card(s,l,1.5,3.85,2.5,st,GREEN)
    txt(s,l,2.35,3.85,0.5,[(t,20,INK,True)],align=PP_ALIGN.CENTER)
    txt(s,l+0.25,3.0,3.35,1.0,[(desc,14,INK2,False)],align=PP_ALIGN.CENTER)
box(s,0.5,4.7,12.3,1.3,fill=LIGHT,line=GREEN_D,line_w=2)
txt(s,0.5,4.85,12.3,1.1,[("AIで速く。人で外さない。データで伸ばす。",30,GREEN_D,True),
  ("“続ける”を超えて、“伸びる”をBtoBへ。",18,INK2,False)],align=PP_ALIGN.CENTER)

prs.save(OUT); print("saved",OUT)
